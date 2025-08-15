import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT


def add_title_and_bullets_slide(prs: Presentation, title: str, bullets: list[str]):
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)

    title_placeholder = slide.shapes.title
    content_placeholder = slide.placeholders[1]

    title_placeholder.text = title
    text_frame = content_placeholder.text_frame

    # Clear default paragraph
    text_frame.clear()

    for idx, bullet in enumerate(bullets):
        p = text_frame.add_paragraph() if idx > 0 else text_frame.paragraphs[0]
        p.text = bullet
        p.level = 0
        p.space_after = Pt(6)
        p.font.size = Pt(20)


def build_presentation() -> Presentation:
    prs = Presentation()

    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "PhishGuard AI"
    slide.placeholders[1].text = "Phishing URL Detection System\nFlask + React + scikit-learn + SQLite"

    # Slides content
    slides_data = [
        (
            "What is Web Phishing?",
            [
                "Fraudulent attempt to obtain sensitive data by impersonating trusted websites",
                "Targets: credentials, banking info, personal data",
                "Tactics: deceptive URLs, spoofed sites, urgency/fear messaging",
            ],
        ),
        (
            "Prerequisites",
            [
                "Software: Python 3.10+, Node.js 16+, npm",
                "Python: Flask, SQLAlchemy, Flask-Login, scikit-learn, BeautifulSoup",
                "Frontend: React 18, Vite, Tailwind CSS, Axios",
                "Database: SQLite (development)",
                "Optional: .env with SECRET_KEY, DATABASE_URI",
            ],
        ),
        (
            "Types of Phishing",
            [
                "Email phishing: mass campaigns with malicious links",
                "Spear phishing: targeted, personalized messages",
                "Whaling: executives/C‑suite targeting",
                "Smishing/Vishing: SMS/voice-based scams",
                "Clone/Website spoofing: look‑alike pages and URLs",
                "Pharming: redirects via DNS/host manipulation",
                "Angler: social media support impersonation",
            ],
        ),
        (
            "System Architecture",
            [
                "Frontend (React): URL input, auth, dashboards",
                "Backend (Flask): routes, auth, scanning API, history",
                "ML Engine: RandomForestClassifier + StandardScaler",
                "Database: SQLite via SQLAlchemy (User, ScanHistory)",
                "APIs: /api/scan, /api/login, /api/register, /api/user/profile",
            ],
        ),
        (
            "Algorithms Used",
            [
                "Model: Random Forest (200 trees, depth 15, class_weight=balanced)",
                "Preprocessing: feature scaling with StandardScaler",
                "Confidence: class probability as confidence score",
                "Risk scoring: blends confidence with feature flags",
                "Persistence: joblib saves/loads model and scaler",
            ],
        ),
        (
            "AI Features Engineered",
            [
                "Structure: lengths (URL/domain/path), subdomains, path depth, query count",
                "Characters: dots, hyphens, underscores, digits, special chars",
                "Security: HTTPS use, IP in domain, custom port",
                "Semantics: suspicious words (verify/login/billing)",
                "Entropy: URL/domain/path/query/fragment Shannon entropy",
                "Heuristics: shorteners, suspicious TLDs, redirect/click/track",
            ],
        ),
        (
            "Application of the Algorithm",
            [
                "Real-time classification via /api/scan (is_phishing + confidence)",
                "Risk levels mapped in backend: Low / Medium / High",
                "Explainability: returns feature snapshot + text sample",
                "History logging: URL, result, score, features per user",
                "Fail-safe: safe defaults on errors",
            ],
        ),
        (
            "Frameworks Used",
            [
                "Backend: Flask, Flask-Login, SQLAlchemy",
                "ML: scikit-learn (RandomForest, StandardScaler), joblib",
                "Scraping: BeautifulSoup",
                "Frontend: React, React Router, Tailwind CSS, Vite",
                "HTTP: Axios",
            ],
        ),
        (
            "Technologies Used",
            [
                "Languages: Python, JavaScript",
                "Database: SQLite (dev-ready, swappable)",
                "Build/Tooling: Vite, npm",
                "Auth: session-based + simple bearer token APIs",
                "Deployment-ready: Gunicorn backend, static frontend build",
            ],
        ),
        (
            "Application of the Project",
            [
                "End-users: validate suspicious links before clicking",
                "Security teams: lightweight triage + history",
                "Education: Info page for awareness and tips",
                "Admin: user management, total scans",
                "Auditability: per-user scan history & timestamps",
            ],
        ),
        (
            "How It Works (User Flow)",
            [
                "Input: user submits a URL",
                "Feature extraction: URL parsed; features computed",
                "Model inference: Random Forest predicts + probability",
                "Risk assessment: confidence → risk level",
                "Result: React UI shows risk, confidence, indicators",
            ],
        ),
        (
            "Risk Classification (Defaults)",
            [
                "Low: confidence > 85% (safe)",
                "Medium: 50–85%",
                "High: < 50% (likely phishing)",
            ],
        ),
        (
            "Data & Training (Project Context)",
            [
                "Training data: synthetic, feature-driven generation",
                "Goal: demonstrate classic feature-based URL detection",
                "Artifacts: phishing_model.pkl, phishing_scaler.pkl via joblib",
            ],
        ),
        (
            "Security & Limitations",
            [
                "Never trust a single signal; combine with user education",
                "Features-based approach; sophisticated evasion is possible",
                "Mitigation: continuous feature updates and retraining",
            ],
        ),
        (
            "Key Files (Reference)",
            [
                "Backend: app.py, phishing_detector.py",
                "Templates: templates/*.html",
                "Frontend: src/App.jsx, src/pages/*, src/components/Navbar.jsx",
                "Model artifacts: phishing_model.pkl, phishing_scaler.pkl",
            ],
        ),
        (
            "Future Enhancements",
            [
                "Data: real labeled datasets for training/eval",
                "Modeling: gradient boosting/transformers for URL text",
                "Explainability: feature importances, SHAP",
                "Integrations: browser extension, threat intel feeds",
            ],
        ),
    ]

    for title, bullets in slides_data:
        add_title_and_bullets_slide(prs, title, bullets)

    return prs


def main():
    prs = build_presentation()
    output_name = "PhishGuard_AI_Presentation.pptx"
    prs.save(output_name)
    print(f"Saved presentation to: {os.path.abspath(output_name)}")


if __name__ == "__main__":
    main()


